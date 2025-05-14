from pptx import Presentation

# Create a new PowerPoint presentation object
ppt = Presentation()

# Slide 1: Title Slide
slide_1 = ppt.slides.add_slide(ppt.slide_layouts[0])
title_1 = slide_1.shapes.title
title_1.text = "半月学习汇报"
subtitle_1 = slide_1.placeholders[1]
subtitle_1.text = "学习成果及下一阶段目标"

# Slide 2: 学习成果 - Python Cookbook
slide_2 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "学习成果 - Python Cookbook"
content_2 = slide_2.shapes.placeholders[1].text_frame
content_2.text = "• 学习 collections.Counter 类"
p = content_2.add_paragraph()
p.text = "• 对比 operator.itemgetter 和 attrgetter 函数"
p = content_2.add_paragraph()
p.text = "• 使用 itertools.groupby 进行数据分组"

# Slide 3: 学习成果 - 吴恩达理论学习
slide_3 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "学习成果 - 吴恩达理论学习"
content_3 = slide_3.shapes.placeholders[1].text_frame
content_3.text = "• 从 Logistic 回归、损失函数、梯度下降算法入手"
p = content_3.add_paragraph()
p.text = "• 对比 for 循环与向量化代码"
p = content_3.add_paragraph()
p.text = "• 研究神经网络正反向传播的数学与向量化表达"

# Slide 4: 学习成果 - 环境配置
slide_4 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "学习成果 - 环境配置"
content_4 = slide_4.shapes.placeholders[1].text_frame
content_4.text = "• 重新配置 Anaconda"
p = content_4.add_paragraph()
p.text = "• 解决 Jupyter Notebook 无法打开工作目录问题"

# Slide 5: 遇到的问题
slide_5 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "遇到的问题"
content_5 = slide_5.shapes.placeholders[1].text_frame
content_5.text = "• 神经网络矩阵表达模糊"
p = content_5.add_paragraph()
p.text = "• 时间管理问题：学习多，实践少"

# Slide 6: 下一阶段目标
slide_6 = ppt.slides.add_slide(ppt.slide_layouts[1])
title_6 = slide_6.shapes.title
title_6.text = "下一阶段目标"
content_6 = slide_6.shapes.placeholders[1].text_frame
content_6.text = "• 深入学习 Python，结合代码实操"
p = content_6.add_paragraph()
p.text = "• 理解 PyTorch 的必要性和功能，再配置环境"
p = content_6.add_paragraph()
p.text = "• 继续吴恩达理论，完成课后习题"

# Save the presentation
ppt.save("D:/Desktop/学习汇报.pptx")
